#!/usr/bin/env python3
"""Parse office/PDF files into a stable JSON contract. Never accesses the business database."""

from __future__ import annotations

import argparse
import csv
import io
import json
import math
import re
import subprocess
from datetime import datetime, timezone
from pathlib import Path


MAX_SUBTITLE_STREAMS = 8
MAX_SUBTITLE_CUES_PER_STREAM = 500
MAX_SUBTITLE_CUE_CHARACTERS = 4_000
MAX_SUBTITLE_BYTES_PER_STREAM = 2 * 1024 * 1024
MAX_SHEET_ROWS_PER_NODE = 200
MAX_SHEET_COLUMNS_PER_NODE = 40
MAX_SHEET_NODE_BYTES = 200 * 1024
MAX_SHEET_CELL_CHARACTERS = 4096
MAX_SHEET_FORMULA_CHARACTERS = 4096
MAX_SHEET_FORMULAS_PER_NODE = 1000
MAX_XLSX_SHEETS = 64
MAX_XLSX_ROWS_PER_SHEET = 50_000
MAX_XLSX_COLUMNS_PER_SHEET = 256
MAX_XLSX_CELLS_PER_SHEET = 1_000_000
MAX_XLSX_TOTAL_CELLS = 2_000_000
MAX_CSV_ROWS = 50_000
MAX_CSV_COLUMNS = 256
MAX_CSV_CELLS = 1_000_000
MAX_SLIDE_NODE_BYTES = 32 * 1024
MAX_IMAGE_PIXELS = 40_000_000
MAX_IMAGE_OCR_LINE_BYTES = 8 * 1024
MAX_IMAGE_OCR_TSV_BYTES = 4 * 1024 * 1024
IMAGE_OCR_TIMEOUT_SECONDS = 30
MAX_PDF_TEXT_BLOCK_BYTES = 32 * 1024
MAX_PDF_TEXT_BLOCKS_PER_PAGE = 500
MAX_PDF_TEXT_BLOCKS_PER_DOCUMENT = 10_000
MAX_PDF_PAGES = 500
MAX_DOCX_PARAGRAPHS = 10_000
MAX_DOCX_NODES = 10_000
MAX_DOCX_PARAGRAPH_BYTES = 32 * 1024
MAX_PPTX_SLIDES = 500
MAX_PPTX_SHAPES_PER_SLIDE = 1_000
MAX_PPTX_NODES = 10_000
SRT_TIMESTAMP_PATTERN = re.compile(
    r"^(\d{2}):(\d{2}):(\d{2}),(\d{3})\s+-->\s+"
    r"(\d{2}):(\d{2}):(\d{2}),(\d{3})(?:\s+.*)?$"
)


def pdf_nodes(file_path: Path, version_id: str) -> list[dict]:
    import fitz

    nodes = []
    with fitz.open(str(file_path)) as document:
        if document.page_count > MAX_PDF_PAGES:
            raise SystemExit("PDF_PAGE_LIMIT")
        for page_number, page in enumerate(document, start=1):
            page_rect = page.rect
            page_nodes = 0
            for block in page.get_text("blocks", sort=True):
                if len(block) < 5:
                    continue
                left, top, right, bottom, raw_text = block[:5]
                text = str(raw_text).strip()
                if not text:
                    continue
                bbox = [float(left), float(top), float(right), float(bottom)]
                if (
                    not all(math.isfinite(value) for value in bbox)
                    or bbox[0] < page_rect.x0
                    or bbox[1] < page_rect.y0
                    or bbox[2] <= bbox[0]
                    or bbox[3] <= bbox[1]
                    or bbox[2] > page_rect.x1
                    or bbox[3] > page_rect.y1
                ):
                    continue
                encoded = text.encode("utf-8")
                truncated = len(encoded) > MAX_PDF_TEXT_BLOCK_BYTES
                if truncated:
                    text = encoded[:MAX_PDF_TEXT_BLOCK_BYTES].decode("utf-8", errors="ignore").rstrip() + "…"
                page_nodes += 1
                if page_nodes > MAX_PDF_TEXT_BLOCKS_PER_PAGE or len(nodes) >= MAX_PDF_TEXT_BLOCKS_PER_DOCUMENT:
                    raise SystemExit("PDF_TEXT_BLOCK_LIMIT")
                block_index = page_nodes
                nodes.append({
                    "schemaVersion": 1,
                    "id": f"page-{page_number}-block-{block_index}",
                    "kind": "paragraph",
                    "title": f"第 {page_number} 页 · 文字块 {block_index}",
                    "content": text,
                    "order": len(nodes),
                    "locator": {
                        "type": "pdf",
                        "resourceVersionId": version_id,
                        "page": page_number,
                        "bbox": bbox,
                    },
                    "metadata": {
                        "pageNumber": page_number,
                        "blockIndex": block_index,
                        "coordinateUnit": "pdf_point",
                        "textTruncated": truncated,
                    },
                })
    return nodes


def docx_nodes(file_path: Path, version_id: str) -> list[dict]:
    from docx import Document

    nodes = []
    heading_stack = []
    document = Document(str(file_path))
    for index, paragraph in enumerate(document.paragraphs):
        if index >= MAX_DOCX_PARAGRAPHS:
            raise SystemExit("DOCX_NODE_LIMIT")
        text = paragraph.text.strip()
        if text:
            encoded = text.encode("utf-8")
            truncated = len(encoded) > MAX_DOCX_PARAGRAPH_BYTES
            if truncated:
                text = encoded[:MAX_DOCX_PARAGRAPH_BYTES].decode("utf-8", errors="ignore").rstrip() + "…"
            if len(nodes) >= MAX_DOCX_NODES:
                raise SystemExit("DOCX_NODE_LIMIT")
            node_id = f"paragraph-{index + 1}"
            style_name = paragraph.style.name or ""
            heading_match = re.match(r"^Heading\s+(\d+)$", style_name, re.IGNORECASE)
            kind = "heading" if heading_match else "paragraph"
            parent_id = None
            if heading_match:
                level = int(heading_match.group(1))
                while heading_stack and heading_stack[-1][0] >= level:
                    heading_stack.pop()
                parent_id = heading_stack[-1][1] if heading_stack else None
            else:
                parent_id = heading_stack[-1][1] if heading_stack else None
            nodes.append({
                "schemaVersion": 1,
                "id": node_id,
                "kind": kind,
                "title": text if kind == "heading" else style_name or node_id,
                "content": text,
                **({"parentId": parent_id} if parent_id else {}),
                "order": len(nodes),
                "locator": {"type": "document", "resourceVersionId": version_id, "nodeId": node_id},
                "metadata": {
                    "paragraphIndex": index,
                    "style": style_name,
                    "textTruncated": truncated,
                },
            })
            if heading_match:
                heading_stack.append((level, node_id))
    return nodes


def bounded_slide_text(value: str) -> tuple[str, bool]:
    encoded = value.encode("utf-8")
    if len(encoded) <= MAX_SLIDE_NODE_BYTES:
        return value, False
    truncated = encoded[:MAX_SLIDE_NODE_BYTES].decode("utf-8", errors="ignore").rstrip()
    return f"{truncated}…", True


def shape_text(shape: object) -> str:
    if getattr(shape, "has_table", False):
        return "\n".join(
            "\t".join(cell.text.strip() for cell in row.cells)
            for row in shape.table.rows
        ).strip()
    if getattr(shape, "has_text_frame", False):
        return shape.text.strip()
    return ""


def pptx_nodes(file_path: Path, version_id: str) -> list[dict]:
    from pptx import Presentation

    nodes = []
    presentation = Presentation(str(file_path))
    if len(presentation.slides) > MAX_PPTX_SLIDES:
        raise SystemExit("PPTX_NODE_LIMIT")
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if len(slide.shapes) > MAX_PPTX_SHAPES_PER_SLIDE:
            raise SystemExit("PPTX_NODE_LIMIT")
        for shape in slide.shapes:
            text = shape_text(shape)
            if not text:
                continue
            if len(nodes) >= MAX_PPTX_NODES:
                raise SystemExit("PPTX_NODE_LIMIT")
            bounded_text, truncated = bounded_slide_text(text)
            shape_id = str(shape.shape_id)
            nodes.append({
                "schemaVersion": 1,
                "id": f"slide-{slide_number}-shape-{shape_id}",
                "kind": "slide",
                "title": f"第 {slide_number} 页 · Shape {shape_id}",
                "content": bounded_text,
                "order": len(nodes),
                "locator": {
                    "type": "slide",
                    "resourceVersionId": version_id,
                    "slide": slide_number,
                    "shapeId": shape_id,
                },
                "metadata": {
                    "slideNumber": slide_number,
                    "shapeId": shape_id,
                    "shapeType": str(shape.shape_type),
                    "isTable": bool(getattr(shape, "has_table", False)),
                    "textTruncated": truncated,
                },
            })
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            if len(nodes) >= MAX_PPTX_NODES:
                raise SystemExit("PPTX_NODE_LIMIT")
            bounded_notes, truncated = bounded_slide_text(notes_text)
            nodes.append({
                "schemaVersion": 1,
                "id": f"slide-{slide_number}-notes",
                "kind": "slide",
                "title": f"第 {slide_number} 页 · 备注",
                "content": bounded_notes,
                "order": len(nodes),
                "locator": {"type": "slide", "resourceVersionId": version_id, "slide": slide_number},
                "metadata": {
                    "slideNumber": slide_number,
                    "contentRole": "notes",
                    "textTruncated": truncated,
                },
            })
    return nodes


def bounded_image_ocr_text(value: str) -> tuple[str, bool]:
    encoded = value.encode("utf-8")
    if len(encoded) <= MAX_IMAGE_OCR_LINE_BYTES:
        return value, False
    truncated = encoded[:MAX_IMAGE_OCR_LINE_BYTES].decode("utf-8", errors="ignore").rstrip()
    return f"{truncated}…", True


def image_dimensions(file_path: Path) -> tuple[int, int]:
    from PIL import Image

    with Image.open(file_path) as image:
        if getattr(image, "n_frames", 1) != 1:
            raise SystemExit("IMAGE_MULTIFRAME_UNSUPPORTED")
        width, height = image.size
        if width <= 0 or height <= 0 or width * height > MAX_IMAGE_PIXELS:
            raise SystemExit("IMAGE_DIMENSIONS_INVALID")
        image.verify()
    return width, height


def image_ocr_nodes(file_path: Path, version_id: str, tesseract: str) -> list[dict]:
    width, height = image_dimensions(file_path)
    try:
        completed = subprocess.run(
            [tesseract, str(file_path), "stdout", "--psm", "6", "-l", "chi_sim+eng", "tsv"],
            check=True,
            capture_output=True,
            text=False,
            timeout=IMAGE_OCR_TIMEOUT_SECONDS,
        )
    except subprocess.TimeoutExpired as error:
        raise SystemExit("IMAGE_OCR_TIMEOUT") from error
    except (FileNotFoundError, subprocess.CalledProcessError) as error:
        raise SystemExit("IMAGE_OCR_FAILED") from error
    if len(completed.stdout) > MAX_IMAGE_OCR_TSV_BYTES:
        raise SystemExit("IMAGE_OCR_OUTPUT_LIMIT")
    try:
        rows = list(csv.DictReader(io.StringIO(completed.stdout.decode("utf-8", errors="strict")), delimiter="\t"))
    except UnicodeDecodeError as error:
        raise SystemExit("IMAGE_OCR_OUTPUT_INVALID") from error
    line_words: dict[tuple[str, str, str, str], list[dict]] = {}
    for row in rows:
        text = (row.get("text") or "").strip()
        if not text or row.get("level") != "5":
            continue
        try:
            left = int(row["left"])
            top = int(row["top"])
            word_width = int(row["width"])
            word_height = int(row["height"])
        except (KeyError, TypeError, ValueError):
            continue
        right = left + word_width
        bottom = top + word_height
        if left < 0 or top < 0 or right <= left or bottom <= top or right > width or bottom > height:
            continue
        key = (row.get("block_num", ""), row.get("par_num", ""), row.get("line_num", ""), row.get("page_num", ""))
        line_words.setdefault(key, []).append({"text": text, "left": left, "top": top, "right": right, "bottom": bottom})
    nodes = []
    for words in line_words.values():
        content, truncated = bounded_image_ocr_text(" ".join(word["text"] for word in words))
        left = min(word["left"] for word in words)
        top = min(word["top"] for word in words)
        right = max(word["right"] for word in words)
        bottom = max(word["bottom"] for word in words)
        node_index = len(nodes) + 1
        nodes.append({
            "schemaVersion": 1,
            "id": f"image-ocr-line-{node_index}",
            "kind": "image",
            "title": f"图片文字第 {node_index} 行",
            "content": content,
            "order": len(nodes),
            "locator": {
                "type": "image",
                "resourceVersionId": version_id,
                "bbox": [left, top, right, bottom],
            },
            "metadata": {
                "contentRole": "ocr_line",
                "imageWidth": width,
                "imageHeight": height,
                "textTruncated": truncated,
            },
        })
    return nodes


def excel_column_label(column: int) -> str:
    label = ""
    value = column
    while value > 0:
        value, remainder = divmod(value - 1, 26)
        label = chr(65 + remainder) + label
    return label


def node_range(start_row: int, end_row: int, column_count: int) -> str:
    return f"A{start_row}:{excel_column_label(column_count)}{end_row}"


def text_cell(value: object) -> str:
    text = "" if value is None else str(value)
    return text if len(text) <= MAX_SHEET_CELL_CHARACTERS else f"{text[:MAX_SHEET_CELL_CHARACTERS]}…"


def table_content(rows: list[list[str]]) -> str:
    return "\n".join("\t".join(row) for row in rows)


def table_blocks(rows, supporting_rows=None):
    blocks = []
    current: list[list[str]] = []
    current_formulas: list[list[str]] = []
    current_start = 0
    current_bytes = 0
    formula_iterator = iter(supporting_rows) if supporting_rows is not None else None
    for row_index, source_row in enumerate(rows, start=1):
        row = list(source_row[:MAX_SHEET_COLUMNS_PER_NODE])
        while row and not row[-1].strip():
            row.pop()
        supporting_row = list(next(formula_iterator, [])[:MAX_SHEET_COLUMNS_PER_NODE]) if formula_iterator else []
        while supporting_row and not supporting_row[-1].strip():
            supporting_row.pop()
        width = max(len(row), len(supporting_row))
        if not width:
            if current:
                blocks.append((current_start, current, current_formulas))
                current = []
                current_formulas = []
                current_bytes = 0
            continue
        row.extend([""] * (width - len(row)))
        row_bytes = len("\t".join(row).encode("utf-8")) + 1
        if current and (
            len(current) >= MAX_SHEET_ROWS_PER_NODE
            or current_bytes + row_bytes > MAX_SHEET_NODE_BYTES
        ):
            blocks.append((current_start, current, current_formulas))
            current = []
            current_formulas = []
            current_bytes = 0
        if not current:
            current_start = row_index
        current.append(row)
        current_formulas.append(supporting_row)
        current_bytes += row_bytes
    if current:
        blocks.append((current_start, current, current_formulas))
    return blocks


def table_nodes(sheets, version_id: str, parser: str) -> list[dict]:
    nodes: list[dict] = []
    for sheet_name, sheet_index, values, formulas in sheets:
        for block_index, (start_row, rows, formula_rows) in enumerate(table_blocks(values, formulas), start=1):
            end_row = start_row + len(rows) - 1
            column_count = max(len(row) for row in rows)
            range_value = node_range(start_row, end_row, column_count)
            formula_items = []
            formula_summary_truncated = False
            if formulas is not None:
                for row_offset, formula_row in enumerate(formula_rows, start=start_row):
                    for column_index, formula in enumerate(formula_row[:column_count], start=1):
                        if isinstance(formula, str) and formula.startswith("="):
                            if len(formula_items) >= MAX_SHEET_FORMULAS_PER_NODE:
                                formula_summary_truncated = True
                                continue
                            if len(formula) > MAX_SHEET_FORMULA_CHARACTERS:
                                formula_summary_truncated = True
                                formula = f"{formula[:MAX_SHEET_FORMULA_CHARACTERS]}…"
                            formula_items.append({
                                "cell": f"{excel_column_label(column_index)}{row_offset}",
                                "formula": formula,
                            })
            nodes.append({
                "schemaVersion": 1,
                "id": f"sheet-{sheet_index}-range-{block_index}",
                "kind": "table",
                "title": f"{sheet_name} · {range_value}",
                "content": table_content(rows),
                "order": len(nodes),
                "locator": {
                    "type": "sheet",
                    "resourceVersionId": version_id,
                    "sheet": sheet_name,
                    "range": range_value,
                },
                "metadata": {
                    "sheetIndex": sheet_index,
                    "sheetName": sheet_name,
                    "range": range_value,
                    "rowStart": start_row,
                    "rowEnd": end_row,
                    "columnCount": column_count,
                    "parser": parser,
                    "formulas": formula_items,
                    "formulaSummaryTruncated": formula_summary_truncated,
                },
            })
    return nodes


def csv_nodes(file_path: Path, version_id: str) -> list[dict]:
    with file_path.open("r", encoding="utf-8", newline="") as handle:
        sample = handle.read(8192)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            dialect = csv.excel
        handle.seek(0)

        def rows():
            total_cells = 0
            for row_index, source_row in enumerate(csv.reader(handle, dialect), start=1):
                cell_count = len(source_row)
                total_cells += cell_count
                if (
                    row_index > MAX_CSV_ROWS
                    or cell_count > MAX_CSV_COLUMNS
                    or total_cells > MAX_CSV_CELLS
                ):
                    raise SystemExit("CSV_DIMENSION_LIMIT")
                yield [text_cell(value) for value in source_row]

        nodes = table_nodes([("CSV", 1, rows(), None)], version_id, "python-csv")
    for node in nodes:
        node["metadata"]["delimiter"] = dialect.delimiter
    return nodes


def xlsx_nodes(file_path: Path, version_id: str) -> list[dict]:
    from openpyxl import load_workbook

    values_workbook = load_workbook(file_path, read_only=True, data_only=True)
    formulas_workbook = load_workbook(file_path, read_only=True, data_only=False)
    try:
        if len(values_workbook.worksheets) > MAX_XLSX_SHEETS:
            raise SystemExit("XLSX_DIMENSION_LIMIT")
        sheets = []
        total_cells = 0
        for sheet_index, value_sheet in enumerate(values_workbook.worksheets, start=1):
            formula_sheet = formulas_workbook[value_sheet.title]
            row_count = max(value_sheet.max_row, formula_sheet.max_row)
            column_count = max(value_sheet.max_column, formula_sheet.max_column)
            cell_count = row_count * column_count
            if (
                row_count > MAX_XLSX_ROWS_PER_SHEET
                or column_count > MAX_XLSX_COLUMNS_PER_SHEET
                or cell_count > MAX_XLSX_CELLS_PER_SHEET
                or total_cells + cell_count > MAX_XLSX_TOTAL_CELLS
            ):
                raise SystemExit("XLSX_DIMENSION_LIMIT")
            total_cells += cell_count
            values = ([text_cell(value) for value in row] for row in value_sheet.iter_rows(values_only=True))
            formulas = ([text_cell(value) for value in row] for row in formula_sheet.iter_rows(values_only=True))
            sheets.append((value_sheet.title, sheet_index, values, formulas))
        return table_nodes(sheets, version_id, "python-xlsx")
    finally:
        values_workbook.close()
        formulas_workbook.close()


def srt_timestamp_to_ms(parts: tuple[str, str, str, str]) -> int:
    hours, minutes, seconds, milliseconds = (int(value) for value in parts)
    return (((hours * 60 + minutes) * 60) + seconds) * 1000 + milliseconds


def parse_srt_cues(value: str, duration_ms: int) -> list[tuple[int, int, str]] | None:
    cues: list[tuple[int, int, str]] = []
    for raw_block in re.split(r"\r?\n\s*\r?\n", value.strip()):
        lines = [line.rstrip() for line in raw_block.splitlines()]
        if not lines:
            continue
        if lines[0].strip().isdigit():
            lines = lines[1:]
        if not lines:
            return None
        timestamp = SRT_TIMESTAMP_PATTERN.match(lines[0].strip())
        if not timestamp:
            return None
        start_ms = srt_timestamp_to_ms(timestamp.groups()[0:4])
        end_ms = srt_timestamp_to_ms(timestamp.groups()[4:8])
        content = "\n".join(line.strip() for line in lines[1:] if line.strip()).strip()
        if (
            not content
            or len(content) > MAX_SUBTITLE_CUE_CHARACTERS
            or start_ms < 0
            or end_ms <= start_ms
            or end_ms > duration_ms
        ):
            return None
        cues.append((start_ms, end_ms, content))
        if len(cues) > MAX_SUBTITLE_CUES_PER_STREAM:
            return None
    return cues or None


def extract_embedded_subtitle_nodes(
    file_path: Path,
    version_id: str,
    duration_ms: int,
    subtitle_streams: list[dict],
    ffmpeg_path: str,
) -> tuple[list[dict], list[dict]]:
    nodes: list[dict] = []
    results: list[dict] = []
    for stream in subtitle_streams[:MAX_SUBTITLE_STREAMS]:
        stream_index = stream.get("index")
        if not isinstance(stream_index, int):
            continue
        try:
            completed = subprocess.run(
                [
                    ffmpeg_path,
                    "-nostdin",
                    "-v",
                    "error",
                    "-i",
                    str(file_path),
                    "-map",
                    f"0:{stream_index}",
                    "-f",
                    "srt",
                    "-",
                ],
                check=True,
                capture_output=True,
                timeout=30,
            )
        except (FileNotFoundError, subprocess.TimeoutExpired, subprocess.CalledProcessError):
            results.append({"index": stream_index, "status": "unavailable"})
            continue
        if len(completed.stdout) > MAX_SUBTITLE_BYTES_PER_STREAM:
            results.append({"index": stream_index, "status": "too_large"})
            continue
        try:
            text = completed.stdout.decode("utf-8", errors="strict")
        except UnicodeDecodeError:
            results.append({"index": stream_index, "status": "invalid_encoding"})
            continue
        cues = parse_srt_cues(text, duration_ms)
        if not cues:
            results.append({"index": stream_index, "status": "unavailable"})
            continue
        language = stream.get("tags", {}).get("language")
        codec = stream.get("codec_name")
        for cue_index, (start_ms, end_ms, content) in enumerate(cues, start=1):
            nodes.append({
                "schemaVersion": 1,
                "id": f"subtitle-{stream_index}-{cue_index}",
                "kind": "transcript",
                "title": "内嵌字幕",
                "content": content,
                "order": 1 + len(nodes),
                "locator": {
                    "type": "video",
                    "resourceVersionId": version_id,
                    "startMs": start_ms,
                    "endMs": end_ms,
                },
                "metadata": {
                    "source": "embedded_subtitle",
                    "streamIndex": stream_index,
                    "codec": codec,
                    "language": str(language)[:32] if language else None,
                },
            })
        results.append({"index": stream_index, "status": "extracted", "cueCount": len(cues)})
    for stream in subtitle_streams[MAX_SUBTITLE_STREAMS:]:
        if isinstance(stream.get("index"), int):
            results.append({"index": stream["index"], "status": "stream_limit"})
    return nodes, results


def media_probe_nodes(
    file_path: Path,
    mime_type: str,
    version_id: str,
    ffprobe_path: str,
    ffmpeg_path: str,
) -> list[dict]:
    """Read local container metadata only; this is not transcription or video understanding."""
    command = [
        ffprobe_path,
        "-v",
        "error",
        "-show_entries",
        "format=format_name,duration,bit_rate:stream=index,codec_type,codec_name,channels,sample_rate,width,height,disposition:stream_tags=language,title",
        "-of",
        "json",
        str(file_path),
    ]
    try:
        completed = subprocess.run(
            command,
            check=True,
            capture_output=True,
            text=True,
            timeout=30,
        )
    except FileNotFoundError as error:
        raise SystemExit("MEDIA_PROBE_UNAVAILABLE") from error
    except subprocess.TimeoutExpired as error:
        raise SystemExit("MEDIA_PROBE_TIMEOUT") from error
    except subprocess.CalledProcessError as error:
        raise SystemExit("MEDIA_PROBE_FAILED") from error

    try:
        probe = json.loads(completed.stdout)
        duration_ms = round(float(probe["format"]["duration"]) * 1000)
    except (KeyError, TypeError, ValueError, json.JSONDecodeError) as error:
        raise SystemExit("MEDIA_PROBE_DURATION_INVALID") from error
    if duration_ms <= 0:
        raise SystemExit("MEDIA_PROBE_DURATION_INVALID")

    streams = probe.get("streams", [])
    audio_streams = [stream for stream in streams if stream.get("codec_type") == "audio"]
    video_streams = [stream for stream in streams if stream.get("codec_type") == "video"]
    subtitle_streams = [stream for stream in streams if stream.get("codec_type") == "subtitle"]
    is_video = mime_type.startswith("video/")
    if (is_video and not video_streams) or (not is_video and not audio_streams):
        raise SystemExit("MEDIA_PROBE_STREAM_MISSING")

    subtitle_nodes: list[dict] = []
    subtitle_extraction: list[dict] = []
    if is_video and subtitle_streams:
        subtitle_nodes, subtitle_extraction = extract_embedded_subtitle_nodes(
            file_path, version_id, duration_ms, subtitle_streams, ffmpeg_path
        )

    summary = {
        "kind": "video" if is_video else "audio",
        "durationMs": duration_ms,
        "container": probe.get("format", {}).get("format_name", "unknown"),
        "audioStreams": [
            {
                "index": stream.get("index"),
                "codec": stream.get("codec_name"),
                "channels": stream.get("channels"),
                "sampleRate": stream.get("sample_rate"),
                "language": str(stream.get("tags", {}).get("language", ""))[:32] or None,
                "default": bool(stream.get("disposition", {}).get("default", 0)),
            }
            for stream in audio_streams
        ],
        "videoStreams": [
            {
                "codec": stream.get("codec_name"),
                "width": stream.get("width"),
                "height": stream.get("height"),
            }
            for stream in video_streams
        ],
        "subtitleStreams": [
            {
                "index": stream.get("index"),
                "codec": stream.get("codec_name"),
                "language": str(stream.get("tags", {}).get("language", ""))[:32] or None,
                "title": str(stream.get("tags", {}).get("title", ""))[:120] or None,
                "default": bool(stream.get("disposition", {}).get("default", 0)),
            }
            for stream in subtitle_streams
        ] if is_video else [],
        "subtitleExtraction": subtitle_extraction if is_video else [],
        "bitRate": probe.get("format", {}).get("bit_rate"),
    }
    locator = {
        "type": "video" if is_video else "audio",
        "resourceVersionId": version_id,
        "startMs": 0,
        "endMs": duration_ms,
    }
    return [{
        "schemaVersion": 1,
        "id": "media-metadata-1",
        "kind": "paragraph",
        "title": "媒体元数据",
        "content": "媒体元数据已提取；尚未生成转写或画面理解。",
        "order": 0,
        "locator": locator,
        "metadata": summary,
    }, *subtitle_nodes]


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    parser.add_argument("--mime", required=True)
    parser.add_argument("--version-id", required=True)
    parser.add_argument("--ffprobe", default="ffprobe")
    parser.add_argument("--ffmpeg", default="ffmpeg")
    parser.add_argument("--tesseract", default="tesseract")
    args = parser.parse_args()
    file_path = Path(args.input).resolve(strict=True)
    handlers = {
        "text/csv": csv_nodes,
        "application/pdf": pdf_nodes,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": docx_nodes,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": pptx_nodes,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": xlsx_nodes,
    }
    image_mime_types = {"image/png", "image/jpeg", "image/webp"}
    media_mime_types = {
        "audio/mpeg",
        "audio/wav",
        "audio/mp4",
        "audio/x-m4a",
        "video/mp4",
        "video/webm",
        "video/quicktime",
    }
    if args.mime in media_mime_types:
        nodes = media_probe_nodes(file_path, args.mime, args.version_id, args.ffprobe, args.ffmpeg)
        parser_id = "wknowledge-python-media-probe"
    elif args.mime in image_mime_types:
        nodes = image_ocr_nodes(file_path, args.version_id, args.tesseract)
        parser_id = "wknowledge-python-image-ocr"
    else:
        handler = handlers.get(args.mime)
        if handler is None:
            raise SystemExit(f"PARSER_MIME_UNSUPPORTED: {args.mime}")
        nodes = handler(file_path, args.version_id)
        parser_id = "wknowledge-python-document"
    if not nodes:
        raise SystemExit("PARSER_EMPTY_RESULT")
    generated_at = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")
    print(json.dumps({
        "document": {
            "schemaVersion": 1,
            "resourceVersionId": args.version_id,
            "nodes": nodes,
        },
        "manifest": {
            "schemaVersion": 1,
            "parserId": parser_id,
            "parserVersion": "1.0.0",
            "runtime": "python",
            "mimeType": args.mime,
            "resourceVersionId": args.version_id,
            "generatedAt": generated_at,
        },
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
